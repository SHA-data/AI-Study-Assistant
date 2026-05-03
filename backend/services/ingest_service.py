import os
import requests
from typing import List
from bs4 import BeautifulSoup
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document as LCDocument
from config import config
from database import SessionLocal
from models import Document as DBDocument

from services.vector_store import get_chroma_store, embedder

def ingest_pdf(file_path: str, namespace: str, title: str = None) -> DBDocument:
    """
    Loads a PDF, chunks it, embeds it into ChromaDB, and records it in SQLite.
    """
    if not title:
        title = os.path.basename(file_path)
    
    loader = PyPDFLoader(file_path)
    documents = loader.load()
    
    return _chunk_and_embed(documents, namespace, title, config.DOC_TYPE_PDF)

def ingest_pptx(file_path: str, namespace: str, title: str = None) -> DBDocument:
    """
    Loads a PPTX, extracts text, chunks it, embeds it into ChromaDB, and records it in SQLite.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required to parse PPTX files.")

    if not title:
        title = os.path.basename(file_path)
    
    prs = Presentation(file_path)
    text_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        if slide_text:
            text_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
    
    full_text = "\n\n".join(text_content)
    documents = [LCDocument(page_content=full_text, metadata={"source": title})]
    
    return _chunk_and_embed(documents, namespace, title, config.DOC_TYPE_PPTX)

def ingest_text(text: str, namespace: str, title: str) -> DBDocument:
    """
    Wraps text in a LangChain Document and processes it.
    """
    documents = [LCDocument(page_content=text, metadata={"source": title})]
    return _chunk_and_embed(documents, namespace, title, config.DOC_TYPE_TEXT)

def ingest_url(url: str, namespace: str, title: str) -> DBDocument:
    """
    Scrapes a URL, saves to uploads, and feeds into the embedding pipeline.
    """
    response = requests.get(url)
    response.raise_for_status()
    
    soup = BeautifulSoup(response.text, 'html.parser')
    
    # Extract p, h1, h2, h3 tags
    tags = soup.find_all(['p', 'h1', 'h2', 'h3'])
    text_content = "\n\n".join([tag.get_text().strip() for tag in tags if tag.get_text().strip()])
    
    # Save scraped text to uploads
    filename = f"scraped_{title.replace(' ', '_')}.txt"
    upload_path = os.path.join(config.UPLOAD_DIR, filename)
    with open(upload_path, "w", encoding="utf-8") as f:
        f.write(text_content)
    
    documents = [LCDocument(page_content=text_content, metadata={"source": url})]
    return _chunk_and_embed(documents, namespace, title, config.DOC_TYPE_URL, source_url=url)

def _chunk_and_embed(
    documents: List[LCDocument], 
    namespace: str, 
    title: str, 
    doc_type: str, 
    source_url: str = None
) -> DBDocument:
    """
    Splits documents into chunks, adds them to ChromaDB, and creates a database record.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100
    )
    chunks = text_splitter.split_documents(documents)
    
    # Add to ChromaDB
    vector_store = get_chroma_store(namespace)
    vector_store.add_documents(chunks)
    
    # Create DB record
    with SessionLocal() as session:
        db_doc = DBDocument(
            title=title,
            doc_type=doc_type,
            namespace=namespace,
            chunk_count=len(chunks),
            source_url=source_url
        )
        session.add(db_doc)
        session.commit()
        session.refresh(db_doc)
        return db_doc
